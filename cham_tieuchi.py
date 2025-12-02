# cham_tieuchi.py
import os, json, zipfile, tempfile, shutil, re, unicodedata
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ===================== HÀM TIỆN ÍCH =====================
def normalize_text_no_diacritics(s):
    if not isinstance(s, str):
        return ""
    s = s.lower()
    s = unicodedata.normalize('NFD', s)
    s = ''.join(ch for ch in s if unicodedata.category(ch) != 'Mn')
    return s

def pretty_name_from_filename(filename):
    base = os.path.splitext(os.path.basename(filename))[0]
    s = base.replace("_", " ").replace("-", " ")
    s = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', s)
    s = re.sub(r'(\d+)', r' \1 ', s)
    parts = [p for p in s.split() if p.strip()]
    parts = [p.capitalize() for p in parts]
    return " ".join(parts)

def find_criteria_file(subject_prefix, grade, criteria_folder="criteria"):
    patterns = [
        f"{subject_prefix}{grade}.json",
        f"{subject_prefix}_khoi{grade}.json",
        f"{subject_prefix}_khoi_{grade}.json",
        f"{subject_prefix}-khoi{grade}.json",
    ]
    for p in patterns:
        full = os.path.join(criteria_folder, p)
        if os.path.exists(full):
            return full
    if os.path.isdir(criteria_folder):
        for fn in os.listdir(criteria_folder):
            lower = fn.lower()
            if fn.lower().startswith(subject_prefix.lower()) and str(grade) in lower and fn.lower().endswith(".json"):
                return os.path.join(criteria_folder, fn)
    return None

def load_criteria(subject, grade, criteria_folder="criteria"):
    """
    subject: 'word'/'ppt'/'scratch' or longer names
    grade: int
    returns dict with key "tieu_chi": list of {'mo_ta','diem', optional 'key','value'}
    """
    if isinstance(grade, str) and grade.isdigit():
        grade = int(grade)
    s = subject.lower()
    if s in ("powerpoint", "ppt", "pptx"):
        pref = "ppt"
    elif s in ("word", "docx", "doc"):
        pref = "word"
    elif s in ("scratch", "sb3"):
        pref = "scratch"
    else:
        pref = s

    path = find_criteria_file(pref, grade, criteria_folder)
    if not path:
        return None
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict) and "tieu_chi" in data:
            for it in data["tieu_chi"]:
                if "diem" in it:
                    try:
                        it["diem"] = float(it["diem"])
                    except:
                        it["diem"] = 0.0
                else:
                    it["diem"] = 0.0
            return data
        elif isinstance(data, list):
            return {"tieu_chi": data}
    except Exception:
        return None
    return None

# ===================== WORD =====================
def grade_word(file_path, criteria):
    try:
        doc = Document(file_path)
    except Exception as e:
        return None, [f"Lỗi đọc file Word: {e}"]

    text = normalize_text_no_diacritics("\n".join([p.text for p in doc.paragraphs]))
    items = criteria.get("tieu_chi", [])
    total_awarded = 0.0
    notes = []

    for it in items:
        desc = it.get("mo_ta", "").strip()
        key = (it.get("key") or "").lower()
        diem = float(it.get("diem", 0) or 0)
        ok = False

        # Các kiểm tra phổ biến
        if key == "has_title":
            ok = any(len(p.text.strip()) > 0 for p in doc.paragraphs[:2])
        elif key == "has_name":
            ok = any(k in text for k in ["ho ten", "ten hoc sinh", "ho va ten"])
        elif key == "has_image":
            ok = bool(getattr(doc, "inline_shapes", None) and len(doc.inline_shapes) > 0)
        elif key == "format_text":
            ok = any(run.bold for p in doc.paragraphs for run in p.runs if getattr(run, "bold", False))
        elif key.startswith("contains:"):
            terms = key.split(":",1)[1].split("|")
            ok = any(normalize_text_no_diacritics(t).strip() in text for t in terms if t.strip())
        else:
            # fallback: tìm từ khoá trong mô tả tiêu chí
            words = re.findall(r"\w+", normalize_text_no_diacritics(desc))
            ok = any(w for w in words if len(w) > 1 and w in text)

        if ok:
            total_awarded += diem
            notes.append(f"✅ {desc} ({diem}đ)")
        else:
            notes.append(f"❌ {desc} (0đ)")

    total_awarded = round(total_awarded, 2)
    return total_awarded, notes

# ===================== POWERPOINT =====================
def grade_ppt(file_path, criteria):
    try:
        prs = Presentation(file_path)
        slides = prs.slides
    except Exception as e:
        return None, [f"Lỗi đọc file PowerPoint: {e}"]

    items = criteria.get("tieu_chi", [])
    total_awarded = 0.0
    notes = []
    num_slides = len(slides)

    def slide_has_picture(slide):
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    return True
                xml_str = getattr(shape, "_element", None)
                if xml_str is not None:
                    xml_str = shape._element.xml.lower()
                    if any(tag in xml_str for tag in ["p:pic", "a:blip", "a:blipfill", "blipfill", ".jpg", ".png", ".gif", ".jpeg"]):
                        return True
            except Exception:
                continue
        return False

    has_picture_any = any(slide_has_picture(slide) for slide in slides)
    has_transition_any = any("transition" in slide._element.xml for slide in slides)
    ppt_text = " ".join(shape.text for slide in slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    ppt_text_noaccent = normalize_text_no_diacritics(ppt_text)

    for it in items:
        desc = it.get("mo_ta", "").strip()
        key = (it.get("key") or "").lower()
        diem = float(it.get("diem", 0) or 0)
        ok = False

        if key == "min_slides":
            val = int(it.get("value", 1))
            ok = num_slides >= val
        elif key == "has_image":
            ok = has_picture_any
        elif key == "has_transition":
            ok = has_transition_any
        elif key == "title_first":
            if slides:
                first = slides[0]
                ok = any(shape.has_text_frame and shape.text.strip() for shape in first.shapes)
        elif key.startswith("contains:"):
            terms = key.split(":",1)[1].split("|")
            ok = any(normalize_text_no_diacritics(t.strip()) in ppt_text_noaccent for t in terms if t.strip())
        elif key == "any":
            ok = True
        else:
            words = re.findall(r"\w+", normalize_text_no_diacritics(desc))
            ok = any(w in ppt_text_noaccent for w in words if len(w) > 1)

        if ok:
            total_awarded += diem
            notes.append(f"✅ {desc} ({diem}đ)")
        else:
            notes.append(f"❌ {desc} (0đ)")

    total_awarded = round(total_awarded, 2)
    return total_awarded, notes

# ===================== SCRATCH =====================
def analyze_sb3_basic(file_path):
    tempdir = None
    try:
        tempdir = tempfile.mkdtemp(prefix="sb3_")
        with zipfile.ZipFile(file_path, 'r') as z:
            z.extractall(tempdir)
        proj_path = os.path.join(tempdir, "project.json")
        if not os.path.exists(proj_path):
            return None, ["Không tìm thấy project.json trong sb3"]
        with open(proj_path, "r", encoding="utf-8") as f:
            proj = json.load(f)
        targets = proj.get("targets", [])
        flags = {
            "has_loop": False, "has_condition": False, "has_interaction": False,
            "has_variable": False, "has_multiple_sprites": False
        }
        sprite_count = sum(1 for t in targets if not t.get("isStage", False))
        if sprite_count >= 2:
            flags["has_multiple_sprites"] = True
        for t in targets:
            if "variables" in t and len(t.get("variables", [])) > 0:
                flags["has_variable"] = True
            blocks = t.get("blocks", {})
            for _, block in blocks.items():
                opcode = (block.get("opcode") or "").lower()
                if any(k in opcode for k in ["control_repeat","control_forever","control_repeat_until"]):
                    flags["has_loop"] = True
                if any(k in opcode for k in ["control_if","control_if_else"]):
                    flags["has_condition"] = True
                if any(k in opcode for k in ["sensing_keypressed","sensing_touchingobject","event_whenthisspriteclicked",
                                             "event_whenflagclicked","event_whenbroadcastreceived","sensing_mousedown"]):
                    flags["has_interaction"] = True
                if any(k in opcode for k in ["data_setvariableto","data_changevariableby","data_hidevariable","data_showvariable"]):
                    flags["has_variable"] = True
                if "event_broadcast" in opcode:
                    flags["has_interaction"] = True
        return flags, []
    except Exception as e:
        return None, [f"Lỗi khi phân tích file Scratch: {e}"]
    finally:
        if tempdir and os.path.exists(tempdir):
            shutil.rmtree(tempdir, ignore_errors=True)

def grade_scratch(file_path, criteria):
    flags, err = analyze_sb3_basic(file_path)
    if flags is None:
        return None, err
    items = criteria.get("tieu_chi", [])
    total_awarded = 0.0
    notes = []
    for it in items:
        desc = it.get("mo_ta","")
        key = (it.get("key") or "").lower()
        diem = float(it.get("diem", 0) or 0)
        ok = False
        # map keys commonly used in criteria files
        if key in ("has_loop","has_condition","has_interaction","has_variable","has_multiple_sprites"):
            ok = bool(flags.get(key, False))
        elif key == "any":
            ok = True
        else:
            # fallback: try to match description keywords to flags
            if "vòng lặp" in desc.lower() or "loop" in desc.lower():
                ok = flags.get("has_loop", False)
            elif "điều kiện" in desc.lower() or "logic" in desc.lower():
                ok = flags.get("has_condition", False)
            elif "broadcast" in desc.lower() or "sự kiện" in desc.lower() or "event" in desc.lower():
                ok = flags.get("has_interaction", False)
        if ok:
            total_awarded += diem
            notes.append(f"✅ {desc} ({diem}đ)")
        else:
            notes.append(f"❌ {desc} (0đ)")
    return round(total_awarded, 2), notes
