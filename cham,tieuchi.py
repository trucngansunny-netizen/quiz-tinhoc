# core.py
import os, json, zipfile, tempfile, shutil

from docx import Document
from pptx import Presentation

# --------- Utilities ----------
def pretty_name_from_filename(filename):
    """
    Lấy tên học sinh từ tên file (bỏ phần đuôi), tách camel/chuỗi liền thành nhiều phần, viết hoa đầu.
    Ví dụ: "LeAnhDung.pptx" -> "Le Anh Dung"
    """
    base = os.path.splitext(os.path.basename(filename))[0]
    # thay các dấu gạch dưới, gạch ngang bằng space
    s = base.replace("_", " ").replace("-", " ")
    # chèn khoảng cách trước chữ hoa nối liền (ví dụ LeAnh -> Le Anh)
    import re
    s = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', s)
    # tách các chữ số ra
    s = re.sub(r'(\d+)', r' \1 ', s)
    parts = [p for p in s.split() if p.strip()]
    parts = [p.capitalize() for p in parts]
    return " ".join(parts)

def load_criteria(subject, grade, criteria_folder="criteria"):
    """
    subject: "word", "ppt", "scratch"
    grade: int 3/4/5
    returns dict structure with keys: 'tieu_chi' (list of {mo_ta, diem, opcode?})
    """
    fname = f"{subject}_khoi{grade}.json"
    path = os.path.join(criteria_folder, fname)
    if not os.path.exists(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

# --------- Word grading (example simple checks) ----------
def grade_word(file_path, criteria):
    try:
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs]).lower()
    except Exception as e:
        return None, [f"Lỗi đọc file Word: {e}"]

    total = 0
    notes = []
    for item in criteria.get("tieu_chi", []):
        desc = item.get("mo_ta", "")
        pts = item.get("diem", 0)
        key = item.get("key", "").lower()  # optional key to check
        ok = False
        # các kiểm tra đơn giản theo key
        if key == "has_title":
            ok = any(k in text for k in ["trường tiểu học", "bài tập", "soạn thảo"])
        elif key == "has_name":
            ok = any(k in text for k in ["họ và tên", "họ tên", "tên học sinh", "hs"])
        elif key == "has_image":
            ok = any("graphicData" in p._element.xml for p in doc.paragraphs)
        elif key == "format_text":
            # check bold or font size in range
            bold_text = any(run.bold for p in doc.paragraphs for run in p.runs if run.bold)
            font_size = any(getattr(run.font, "size", None) and 12 <= run.font.size.pt <= 20
                            for p in doc.paragraphs for run in p.runs if getattr(run, "font", None))
            ok = bold_text or font_size
        elif key == "any":
            ok = True
        else:
            # fallback: check keyword in description text
            ok = any(k.lower() in text for k in desc.split())
        if ok:
            total += pts
            notes.append("✅ " + desc)
        else:
            notes.append("❌ " + desc)
    return total, notes

# --------- PPT grading ----------
def grade_ppt(file_path, criteria):
    try:
        prs = Presentation(file_path)
        slides = prs.slides
    except Exception as e:
        return None, [f"Lỗi đọc file PowerPoint: {e}"]
    total = 0
    notes = []
    for item in criteria.get("tieu_chi", []):
        desc = item.get("mo_ta", "")
        pts = item.get("diem", 0)
        key = item.get("key", "").lower()
        ok = False
        if key == "min_slides":
            ok = len(slides) >= item.get("value", 1)
        elif key == "title_first":
            if slides:
                first = slides[0]
                ok = any(shape.has_text_frame and shape.text.strip() for shape in first.shapes)
            else:
                ok = False
        elif key == "has_image":
            ok = any(getattr(shape, "shape_type", None) == 13 for slide in slides for shape in slide.shapes)
        elif key == "has_transition":
            ok = any("transition" in slide._element.xml for slide in slides)
        elif key == "format_text":
            bold_or_colored = False
            for slide in slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if run.font.bold:
                                    bold_or_colored = True; break
                                if run.font.color and hasattr(run.font.color, "rgb") and run.font.color.rgb:
                                    bold_or_colored = True; break
                            except Exception:
                                continue
                        if bold_or_colored: break
                    if bold_or_colored: break
                if bold_or_colored: break
            ok = bold_or_colored
        elif key == "any":
            ok = True
        if ok:
            total += pts; notes.append("✅ " + desc)
        else:
            notes.append("❌ " + desc)
    return total, notes

# --------- Scratch grading ----------
def analyze_sb3_basic(file_path):
    tempdir=None
    try:
        tempdir = tempfile.mkdtemp(prefix="sb3_")
        with zipfile.ZipFile(file_path, 'r') as z:
            z.extractall(tempdir)
        proj_path = os.path.join(tempdir, "project.json")
        if not os.path.exists(proj_path):
            return None, ["Không tìm thấy project.json"]
        with open(proj_path, "r", encoding="utf-8") as f:
            proj = json.load(f)
        targets = proj.get("targets", [])
        # flags
        flags = {
            "has_loop": False,
            "has_condition": False,
            "has_interaction": False,
            "has_variable": False,
            "multiple_sprites_or_animation": False
        }
        sprite_count = sum(1 for t in targets if not t.get("isStage", False))
        if sprite_count >= 2:
            flags["multiple_sprites_or_animation"] = True
        for t in targets:
            costumes = t.get("costumes", [])
            if len(costumes) >= 2:
                flags["multiple_sprites_or_animation"] = True
            if "variables" in t and len(t["variables"]) > 0:
                flags["has_variable"] = True
            blocks = t.get("blocks", {})
            for block_id, block in blocks.items():
                opcode = block.get("opcode", "").lower()
                if any(k in opcode for k in ["control_repeat","control_forever","control_repeat_until"]):
                    flags["has_loop"] = True
                if any(k in opcode for k in ["control_if","control_if_else"]):
                    flags["has_condition"] = True
                if any(k in opcode for k in ["sensing_keypressed","sensing_touchingobject",
                                             "event_whenthisspriteclicked","event_whenstageclicked",
                                             "event_whenflagclicked","event_whenbroadcastreceived",
                                             "sensing_mousedown"]):
                    flags["has_interaction"] = True
                if any(k in opcode for k in ["data_setvariableto","data_changevariableby","data_hidevariable","data_showvariable"]):
                    flags["has_variable"] = True
                if "event_broadcast" in opcode:
                    flags["has_interaction"] = True
        return flags, []
    except Exception as e:
        return None, [f"Lỗi khi phân tích file Scratch: {e}"]
    finally:
        if tempdir and os.path.exists(tempdir):
            try:
                shutil.rmtree(tempdir)
            except Exception:
                pass

def grade_scratch(file_path, criteria):
    flags, err = analyze_sb3_basic(file_path)
    if flags is None:
        return None, err
    total = 0; notes = []
    for item in criteria.get("tieu_chi", []):
        desc = item.get("mo_ta",""); pts = item.get("diem",0); key=item.get("key","").lower()
        ok=False
        if key=="has_loop":
            ok = flags.get("has_loop", False)
        elif key=="has_condition":
            ok = flags.get("has_condition", False)
        elif key=="has_interaction":
            ok = flags.get("has_interaction", False)
        elif key=="has_variable":
            ok = flags.get("has_variable", False)
        elif key=="multiple_sprites_or_animation":
            ok = flags.get("multiple_sprites_or_animation", False)
        elif key=="any":
            ok = True
        if ok:
            total += pts; notes.append("✅ "+desc)
        else:
            notes.append("❌ "+desc)
    return total, notes